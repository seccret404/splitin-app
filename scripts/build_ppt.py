"""Generate a clean, formal, minimal portfolio deck for Edward Tua Panjaitan."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----------------------------------------------------------------
INK = RGBColor(0x22, 0x2A, 0x33)     # near-black slate (text)
MUTED = RGBColor(0x6B, 0x72, 0x80)   # secondary text
ACCENT = RGBColor(0x2F, 0x7E, 0x83)  # deep teal (single accent)
RULE = RGBColor(0xE3, 0xE7, 0xEA)    # hairline
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0xCE, 0xD4, 0xD9)   # large page numbers

SERIF = "Georgia"      # titles — formal/editorial
SANS = "Calibri"       # body

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
LM = Inches(0.95)  # left margin


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return s


def text(s, l, t, w, h, runs, size=16, color=INK, bold=False, font=SANS,
         align=PP_ALIGN.LEFT, italic=False, spacing=1.1, space_after=0, tracking=None):
    """runs: str or list of (text, dict-overrides)."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    if space_after:
        p.space_after = Pt(space_after)
    items = runs if isinstance(runs, list) else [(runs, {})]
    for txt, ov in items:
        r = p.add_run()
        r.text = txt
        f = r.font
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
        f.name = ov.get("font", font)
        f.color.rgb = ov.get("color", color)
        tr = ov.get("tracking", tracking)
        if tr is not None:
            r._r.get_or_add_rPr().set("spc", str(tr))
    return tb


def rule(s, l, t, w, h, color=ACCENT):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def kicker(s, label):
    text(s, LM, Inches(0.85), Inches(8), Inches(0.4),
         label.upper(), size=13, color=ACCENT, bold=True, tracking=2600)
    rule(s, LM, Inches(1.28), Inches(0.62), Pt(3))


def page_no(s, n):
    text(s, SW - Inches(1.3), SH - Inches(0.75), Inches(0.9), Inches(0.4),
         f"{n:02d}", size=12, color=MUTED, align=PP_ALIGN.RIGHT)
    text(s, LM, SH - Inches(0.75), Inches(6), Inches(0.4),
         "Edward Tua Panjaitan", size=10, color=MUTED, tracking=600)


# ===== Slide 1 — Title =======================================================
s = slide()
rule(s, Inches(0), Inches(0), SW, Inches(0.16))  # top accent band
text(s, LM, Inches(2.05), Inches(10), Inches(0.5),
     "SOFTWARE DEVELOPER", size=14, color=ACCENT, bold=True, tracking=3200)
text(s, LM, Inches(2.55), Inches(11.4), Inches(1.8),
     "Edward Tua Panjaitan", size=50, color=INK, bold=True, font=SERIF)
rule(s, LM, Inches(3.95), Inches(0.9), Pt(4))
text(s, LM, Inches(4.2), Inches(10.5), Inches(1.0),
     "Building AI-powered systems that solve real problems.",
     size=20, color=MUTED, font=SERIF, italic=True)
text(s, LM, SH - Inches(1.1), Inches(11), Inches(0.5),
     [("edwardtua25@gmail.com", {}),
      ("      linkedin.com/in/edward-tua-panjaitan", {"color": MUTED})],
     size=13, color=INK, bold=True)

# ===== Slide 2 — About & Skills =============================================
s = slide()
kicker(s, "About")
text(s, LM, Inches(1.55), Inches(11), Inches(0.9),
     "Who I Am", size=34, color=INK, bold=True, font=SERIF)
text(s, LM, Inches(2.5), Inches(11.4), Inches(1.4),
     "Software developer focused on backend systems and AI integration. I build scalable "
     "applications, data pipelines, and intelligent tools — shipping work that is tested, "
     "documented, and built to last.",
     size=16, color=MUTED, spacing=1.3)

text(s, LM, Inches(4.05), Inches(8), Inches(0.4),
     "SKILLS", size=13, color=ACCENT, bold=True, tracking=2600)

groups = [
    ("Languages", "Python · TypeScript · JavaScript"),
    ("Backend", "FastAPI · Node.js · REST APIs"),
    ("AI / ML", "RAG · LangChain · Gemini · OpenAI · Ollama"),
    ("Mobile", "React Native · Expo"),
    ("Tools", "Git · Docker · aiohttp"),
]
y = 4.55
for label, items in groups:
    text(s, LM, Inches(y), Inches(2.4), Inches(0.4), label,
         size=14, color=INK, bold=True)
    text(s, LM + Inches(2.5), Inches(y), Inches(8.6), Inches(0.4), items,
         size=14, color=MUTED)
    y += 0.5
page_no(s, 2)

# ===== Slide 3 — Project 01 =================================================
s = slide()
kicker(s, "Project 01")
text(s, LM, Inches(1.55), Inches(11.4), Inches(0.9),
     "AI Web Crawl & Extraction Platform", size=30, color=INK, bold=True, font=SERIF)
text(s, LM, Inches(2.5), Inches(11.4), Inches(1.2),
     "A full-stack app that crawls any website and uses LLMs to turn messy web content into "
     "clean, structured documents — ready for RAG pipelines.",
     size=16, color=MUTED, spacing=1.3)

text(s, LM, Inches(3.85), Inches(8), Inches(0.4),
     "HIGHLIGHTS", size=13, color=ACCENT, bold=True, tracking=2600)
points = [
    "Async crawler with live progress (server-sent events)",
    "Provider-agnostic LLMs — Gemini, OpenAI, Claude, Ollama",
    "User-defined extraction schema → structured output",
]
y = 4.35
for p in points:
    rule(s, LM, Inches(y + 0.12), Inches(0.18), Pt(3))
    text(s, LM + Inches(0.42), Inches(y), Inches(10.6), Inches(0.5), p,
         size=15.5, color=INK)
    y += 0.62
text(s, LM, SH - Inches(1.15), Inches(11), Inches(0.4),
     "Python · FastAPI · aiohttp · LangExtract", size=13, color=MUTED, bold=True, tracking=800)
page_no(s, 3)

# ===== Slide 4 — Project 02 SplitIn ========================================
s = slide()
kicker(s, "Project 02")
text(s, LM, Inches(1.55), Inches(11.4), Inches(0.9),
     "SplitIn — Bill Splitting for Gen Z", size=30, color=INK, bold=True, font=SERIF)
text(s, LM, Inches(2.5), Inches(11.4), Inches(1.2),
     "A mobile app that takes the chaos out of splitting a bill: enter the menu, mark who "
     "ordered what, and settle up — designed around how people here actually pay and message.",
     size=16, color=MUTED, spacing=1.3)

text(s, LM, Inches(3.85), Inches(8), Inches(0.4),
     "HIGHLIGHTS", size=13, color=ACCENT, bold=True, tracking=2600)
points = [
    "End-to-end flow: Calculate → Share → Collect",
    "Local-first: PB1 tax, QRIS / GoPay / OVO, WhatsApp billing",
    "35 unit tests · type-checked · bilingual (ID / EN)",
]
y = 4.35
for p in points:
    rule(s, LM, Inches(y + 0.12), Inches(0.18), Pt(3))
    text(s, LM + Inches(0.42), Inches(y), Inches(10.6), Inches(0.5), p,
         size=15.5, color=INK)
    y += 0.62
text(s, LM, SH - Inches(1.15), Inches(11), Inches(0.4),
     "React Native · Expo · TypeScript", size=13, color=MUTED, bold=True, tracking=800)
page_no(s, 4)

# ===== Slide 5 — Closing ====================================================
s = slide()
rule(s, Inches(0), SH - Inches(0.16), SW, Inches(0.16))  # bottom accent band
text(s, LM, Inches(2.5), Inches(11.4), Inches(1.3),
     "Let's build something\ntogether.", size=42, color=INK, bold=True, font=SERIF, spacing=1.05)
rule(s, LM, Inches(4.35), Inches(0.9), Pt(4))
text(s, LM, Inches(4.6), Inches(11), Inches(0.6),
     "Open to roles in backend, AI/ML integration, and full-stack development.",
     size=17, color=MUTED, font=SERIF, italic=True)
text(s, LM, Inches(5.5), Inches(11), Inches(0.5),
     [("edwardtua25@gmail.com", {}),
      ("      linkedin.com/in/edward-tua-panjaitan", {"color": MUTED})],
     size=14, color=INK, bold=True)
text(s, LM, SH - Inches(0.95), Inches(11), Inches(0.4),
     "© 2026 Edward Tua Panjaitan · Software Developer Portfolio",
     size=10, color=MUTED, tracking=600)

prs.save("Edward_Portfolio.pptx")
print("Saved Edward_Portfolio.pptx ·", len(prs.slides._sldIdLst), "slides")
